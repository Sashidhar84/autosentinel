"""
AutoSentinel — Automotive Customer Voice Intelligence Platform
Complete Streamlit app. Single file. Deploy on Streamlit Cloud.
Zero cost. Zero hallucination. All signal.
"""

import streamlit as st
import requests
import json
import re
import time
import hashlib
from datetime import datetime, timedelta
from typing import Optional
import os

# ============================================================
# PAGE CONFIG — must be first Streamlit call
# ============================================================

st.set_page_config(
    page_title="AutoSentinel",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ============================================================
# CUSTOM CSS — Dark automotive theme
# ============================================================

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display&family=IBM+Plex+Mono:wght@400;500&family=Sora:wght@300;400;500;600&display=swap');

/* Global */
html, body, [class*="css"] {
    font-family: 'Sora', sans-serif;
    background-color: #0A0A0A;
    color: #F0F0F0;
}
.stApp { background-color: #0A0A0A; }
.block-container { padding-top: 1rem; max-width: 1200px; }

/* Hide default streamlit elements */
#MainMenu, footer, header { visibility: hidden; }

/* Header */
.as-header {
    display: flex; align-items: center; justify-content: space-between;
    padding: 16px 0 16px 0;
    border-bottom: 1px solid #2A2A2A;
    margin-bottom: 32px;
}
.as-logo {
    font-family: 'DM Serif Display', serif;
    font-size: 28px; letter-spacing: 2px;
}
.as-logo span.red { color: #E63946; }
.as-logo span.white { color: #F0F0F0; }
.as-tagline {
    font-family: 'Sora', sans-serif;
    font-size: 12px; color: #666; letter-spacing: 1px;
    text-transform: uppercase;
}

/* Section headers */
.as-section-title {
    font-family: 'DM Serif Display', serif;
    font-size: 22px; color: #F0F0F0;
    margin: 28px 0 12px 0;
}

/* Defect card */
.defect-card {
    background: #141414;
    border: 1px solid #2A2A2A;
    border-radius: 10px;
    padding: 24px;
    margin-bottom: 20px;
}
.defect-rank {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 11px; color: #666;
    text-transform: uppercase; letter-spacing: 1px;
}
.defect-title {
    font-family: 'DM Serif Display', serif;
    font-size: 22px; color: #F0F0F0;
    margin: 8px 0 16px 0;
    line-height: 1.3;
}
.severity-critical {
    background: #E63946; color: #fff;
    padding: 3px 10px; border-radius: 4px;
    font-size: 11px; font-weight: 600;
    letter-spacing: 1px; font-family: 'IBM Plex Mono', monospace;
}
.severity-high {
    background: #F4A261; color: #1a1a1a;
    padding: 3px 10px; border-radius: 4px;
    font-size: 11px; font-weight: 600;
    letter-spacing: 1px; font-family: 'IBM Plex Mono', monospace;
}
.severity-moderate {
    background: #52B788; color: #1a1a1a;
    padding: 3px 10px; border-radius: 4px;
    font-size: 11px; font-weight: 600;
    letter-spacing: 1px; font-family: 'IBM Plex Mono', monospace;
}
.mention-count {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 13px; color: #E63946; font-weight: 600;
}

/* Condition box */
.condition-box {
    background: #0A0A0A;
    border-left: 3px solid #E63946;
    border-radius: 0 6px 6px 0;
    padding: 14px 16px;
    margin: 14px 0;
}
.condition-label {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 10px; color: #E63946;
    text-transform: uppercase; letter-spacing: 2px;
    margin-bottom: 10px;
}
.condition-row {
    display: flex; gap: 8px;
    margin: 5px 0;
    font-size: 13px;
}
.condition-key {
    color: #666; min-width: 120px;
    font-family: 'IBM Plex Mono', monospace;
    font-size: 12px;
}
.condition-val { color: #F0F0F0; }
.condition-null { color: #444; font-style: italic; }

/* Verbatim */
.verbatim-box {
    border-left: 3px solid #2A2A2A;
    padding: 12px 16px;
    margin: 14px 0;
}
.verbatim-label {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 10px; color: #666;
    text-transform: uppercase; letter-spacing: 2px;
    margin-bottom: 8px;
}
.verbatim-quote {
    font-size: 14px; color: #E8E8E8;
    font-style: italic; line-height: 1.8;
}
.verbatim-attr {
    font-size: 11px; color: #555;
    margin-top: 8px; font-family: 'IBM Plex Mono', monospace;
}

/* Source badges */
.source-badge {
    display: inline-block;
    padding: 3px 8px; border-radius: 4px;
    font-size: 11px; font-family: 'IBM Plex Mono', monospace;
    margin: 2px 3px;
    background: #1C1C1C; border: 1px solid #2A2A2A; color: #A0A0A0;
}

/* Special mention / escalation */
.escalation-banner {
    background: #1A0505;
    border: 1px solid #E63946;
    border-radius: 8px;
    padding: 16px 20px;
    margin-bottom: 24px;
}
.escalation-title {
    font-family: 'DM Serif Display', serif;
    font-size: 18px; color: #E63946;
    margin-bottom: 4px;
}
.escalation-card {
    background: #0F0505;
    border-left: 4px solid #E63946;
    border-radius: 0 6px 6px 0;
    padding: 14px 16px;
    margin: 10px 0;
}
.escalation-type {
    background: #E63946; color: #fff;
    padding: 2px 8px; border-radius: 3px;
    font-size: 10px; font-weight: 600;
    letter-spacing: 1px; font-family: 'IBM Plex Mono', monospace;
    display: inline-block; margin-bottom: 8px;
}

/* Summary strip */
.summary-strip {
    background: #141414;
    border: 1px solid #2A2A2A;
    border-radius: 8px;
    padding: 20px 24px;
    margin-bottom: 24px;
}
.big-number {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 48px; color: #E63946;
    line-height: 1; font-weight: 600;
}
.big-number-label {
    font-size: 12px; color: #666;
    text-transform: uppercase; letter-spacing: 1px;
    margin-top: 4px;
}
.car-name-display {
    font-family: 'DM Serif Display', serif;
    font-size: 32px; color: #F0F0F0;
}

/* Progress */
.progress-step {
    display: flex; align-items: center; gap: 12px;
    padding: 8px 0; border-bottom: 1px solid #1A1A1A;
    font-size: 13px;
}
.step-num {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 11px; color: #E63946;
    min-width: 20px;
}

/* Proposal box */
.proposal-box {
    background: #0A0F0F;
    border: 1px solid #2EC4B6;
    border-radius: 8px;
    padding: 20px 24px;
    margin: 8px 0;
}
.proposal-title {
    font-family: 'DM Serif Display', serif;
    font-size: 18px; color: #2EC4B6;
    margin-bottom: 12px;
}

/* Disclaimer */
.disclaimer-pill {
    display: inline-block;
    background: #141414; border: 1px solid #2A2A2A;
    border-radius: 20px; padding: 4px 14px;
    font-size: 11px; color: #666;
    font-family: 'IBM Plex Mono', monospace;
}

/* Warning */
.warning-box {
    background: #0F0A00;
    border: 1px solid #F4A261;
    border-radius: 8px;
    padding: 14px 18px;
    margin: 12px 0;
    font-size: 13px; color: #F4A261;
}

/* Stray streamlit button overrides */
.stButton > button {
    background: #E63946; color: #fff;
    border: none; border-radius: 6px;
    font-family: 'Sora', sans-serif;
    font-weight: 600; font-size: 14px;
    padding: 10px 24px;
    transition: all 0.2s;
}
.stButton > button:hover {
    background: #c1121f; color: #fff;
    border: none;
}
.stSelectbox label, .stTextInput label {
    color: #A0A0A0; font-size: 12px;
    text-transform: uppercase; letter-spacing: 1px;
    font-family: 'IBM Plex Mono', monospace;
}
.stTextInput > div > div > input {
    background: #141414; color: #F0F0F0;
    border: 1px solid #2A2A2A; border-radius: 6px;
}
.stSelectbox > div > div {
    background: #141414; color: #F0F0F0;
    border: 1px solid #2A2A2A; border-radius: 6px;
}
.stTabs [data-baseweb="tab"] {
    color: #666; font-family: 'Sora', sans-serif;
    font-size: 13px;
}
.stTabs [aria-selected="true"] {
    color: #E63946;
    border-bottom: 2px solid #E63946;
}
</style>
""", unsafe_allow_html=True)

# ============================================================
# INDIAN CAR MODELS CONFIG
# ============================================================

INDIAN_CAR_MODELS = {
    "Maruti Suzuki": [
        "Swift", "Baleno", "Brezza", "Ertiga", "XL6", "Fronx", "Jimny",
        "Grand Vitara", "Dzire", "WagonR", "Alto K10", "S-Presso",
        "Celerio", "Ignis", "Ciaz",
    ],
    "Hyundai": [
        "Creta", "Venue", "i20", "Grand i10 Nios", "Aura", "Verna",
        "Alcazar", "Tucson", "Ioniq 5", "Exter",
    ],
    "Tata": [
        "Nexon", "Punch", "Tiago", "Tigor", "Altroz", "Harrier",
        "Safari", "Curvv", "Nexon EV", "Punch EV", "Tiago EV",
    ],
    "Mahindra": [
        "XUV700", "XUV300", "XUV400", "Thar", "Bolero", "Scorpio N",
        "Scorpio Classic", "BE 6e", "XEV 9e", "Bolero Neo",
    ],
    "Honda": ["Amaze", "City", "Elevate", "WR-V"],
    "Toyota": [
        "Innova Crysta", "Innova HyCross", "Fortuner", "Legender",
        "Urban Cruiser Hyryder", "Glanza", "Rumion", "Camry",
    ],
    "Kia": ["Seltos", "Sonet", "Carens", "EV6", "EV9"],
    "MG": ["Hector", "Astor", "Comet EV", "Windsor EV", "Gloster", "ZS EV"],
    "Skoda": ["Slavia", "Kushaq", "Kodiaq", "Superb", "Octavia"],
    "Volkswagen": ["Taigun", "Virtus", "Tiguan", "Passat"],
    "Renault": ["Kwid", "Triber", "Kiger"],
    "Nissan": ["Magnite", "GT-R"],
    "Citroën": ["C3", "C3 Aircross", "C5 Aircross"],
    "Jeep": ["Compass", "Meridian", "Wrangler", "Grand Cherokee"],
    "BMW": ["3 Series", "5 Series", "7 Series", "X1", "X3", "X5", "iX", "i4"],
    "Mercedes-Benz": ["C-Class", "E-Class", "S-Class", "GLA", "GLC", "GLE", "EQB", "EQS"],
    "Audi": ["A4", "A6", "Q3", "Q5", "Q7", "e-tron", "e-tron GT"],
    "Volvo": ["XC40", "XC60", "XC90", "S60", "S90"],
    "Force": ["Gurkha", "Trax"],
    "Isuzu": ["D-Max V-Cross", "MU-X"],
}

# ============================================================
# API HELPERS — All free tier
# ============================================================

def get_secrets():
    """Get API keys from Streamlit secrets or environment variables"""

    def get_key(key_name, fallback=""):
        # Try st.secrets first
        try:
            val = st.secrets[key_name]
            if val and str(val).strip():
                return str(val).strip()
        except Exception:
            pass
        # Try environment variable as fallback
        val = os.environ.get(key_name, fallback)
        return str(val).strip() if val else fallback

    firecrawl = get_key("FIRECRAWL_API_KEY")
    gemini = get_key("GEMINI_API_KEY")
    youtube = get_key("YOUTUBE_API_KEY")

    # Check if critical keys are present
    if not firecrawl or not gemini or not youtube:
        st.error(
            "⚠ API keys not configured. "
            "Go to Streamlit Cloud → Your app → ⋮ menu → Settings → Secrets "
            "and add your keys in TOML format."
        )
        st.code(
            'FIRECRAWL_API_KEY = "your-firecrawl-key"\n'
            'GEMINI_API_KEY = "your-gemini-key"\n'
            'YOUTUBE_API_KEY = "your-youtube-key"\n'
            'REDDIT_CLIENT_ID = "not-applicable"\n'
            'REDDIT_CLIENT_SECRET = "not-applicable"',
            language="toml",
        )
        st.stop()

    return {
        "firecrawl": firecrawl,
        "gemini": gemini,
        "youtube": youtube,
        "reddit_id": get_key("REDDIT_CLIENT_ID", "not-applicable"),
        "reddit_secret": get_key("REDDIT_CLIENT_SECRET", "not-applicable"),
    }

# ============================================================
# DATA FETCHER: YouTube
# ============================================================

def fetch_youtube(api_key: str, brand: str, model: str, year: int,
                  progress_cb=None) -> list:
    comments = []
    queries = [
        f"{brand} {model} {year} problem issue owner",
        f"{brand} {model} {year} defect fault review",
        f"{brand} {model} vibration noise rattle",
        f"{brand} {model} long term review problems",
    ]

    for qi, query in enumerate(queries):
        try:
            search_url = (
                f"https://www.googleapis.com/youtube/v3/search"
                f"?part=snippet&q={requests.utils.quote(query)}"
                f"&type=video&maxResults=8&key={api_key}"
                f"&relevanceLanguage=en&regionCode=IN"
            )
            search_resp = requests.get(search_url, timeout=15)
            search_data = search_resp.json()
            if "items" not in search_data:
                continue

            for video in search_data["items"][:6]:
                video_id = video["id"].get("videoId")
                if not video_id:
                    continue
                video_url = f"https://youtube.com/watch?v={video_id}"
                page_token = None
                page_count = 0

                while page_count < 4:
                    params = {
                        "part": "snippet",
                        "videoId": video_id,
                        "maxResults": 100,
                        "key": api_key,
                    }
                    if page_token:
                        params["pageToken"] = page_token

                    comm_resp = requests.get(
                        "https://www.googleapis.com/youtube/v3/commentThreads",
                        params=params, timeout=15
                    )
                    comm_data = comm_resp.json()
                    if "items" not in comm_data:
                        break

                    for item in comm_data["items"]:
                        c = item["snippet"]["topLevelComment"]["snippet"]
                        if c.get("likeCount", 0) >= 2 or \
                                item["snippet"].get("totalReplyCount", 0) > 0:
                            comments.append({
                                "text": c.get("textDisplay", ""),
                                "username": c.get("authorDisplayName", "YouTube user"),
                                "platform": "youtube",
                                "source_url": video_url,
                                "date": c.get("publishedAt", ""),
                            })

                    page_token = comm_data.get("nextPageToken")
                    if not page_token:
                        break
                    page_count += 1

        except Exception as e:
            pass  # Source fails gracefully

    if progress_cb:
        progress_cb(f"YouTube: {len(comments)} comments collected")
    return comments


# ============================================================
# DATA FETCHER: Team-BHP via Firecrawl
# ============================================================

def fetch_teambhp(firecrawl_key: str, brand: str, model: str,
                  progress_cb=None) -> list:
    posts = []
    model_slug = model.lower().replace(" ", "-")

    urls_to_crawl = [
        f"https://www.team-bhp.com/forum/search.php?q={requests.utils.quote(model)}&type=posts&showposts=1",
        f"https://www.team-bhp.com/forum/technical-stuff/",
        f"https://www.team-bhp.com/forum/car-issues-help/",
    ]

    for url in urls_to_crawl:
        try:
            resp = requests.post(
                "https://api.firecrawl.dev/v1/crawl",
                headers={
                    "Authorization": f"Bearer {firecrawl_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "url": url,
                    "limit": 30,
                    "scrapeOptions": {
                        "formats": ["markdown"],
                        "onlyMainContent": True,
                    },
                },
                timeout=60,
            )
            data = resp.json()
            pages = data.get("data", [])

            for page in pages:
                markdown = page.get("markdown", "")
                if not markdown or model.lower() not in markdown.lower():
                    continue
                blocks = re.split(r'\n---\n|\n_{3,}\n|\n\n\n', markdown)
                for block in blocks:
                    if len(block) > 100 and _has_quality_signal(block):
                        posts.append({
                            "text": block[:2000],
                            "username": _extract_username(block),
                            "platform": "team_bhp",
                            "source_url": page.get("metadata", {}).get(
                                "sourceURL", url
                            ),
                            "date": page.get("metadata", {}).get(
                                "publishedTime", ""
                            ),
                        })
        except Exception:
            pass

    if progress_cb:
        progress_cb(f"Team-BHP: {len(posts)} posts collected")
    return posts


# ============================================================
# DATA FETCHER: Reddit via PRAW API
# ============================================================

def fetch_reddit(client_id: str, client_secret: str,
                 brand: str, model: str, progress_cb=None) -> list:
    posts = []

    # Get access token
    try:
        token_resp = requests.post(
            "https://www.reddit.com/api/v1/access_token",
            auth=(client_id, client_secret),
            data={"grant_type": "client_credentials"},
            headers={"User-Agent": "AutoSentinel/1.0"},
            timeout=15,
        )
        token = token_resp.json().get("access_token")
        if not token:
            return posts
    except Exception:
        return posts

    headers = {
        "Authorization": f"Bearer {token}",
        "User-Agent": "AutoSentinel/1.0",
    }

    subreddits = [
        "CarTalkIndia", "india", "cars", "IndianCars",
        "whatcarshouldibuy", "MechanicAdvice",
    ]
    search_terms = [
        f"{model} problem", f"{model} issue", f"{model} defect",
        f"{model} vibration", f"{model} noise", f"{brand} {model} fault",
    ]

    for sub in subreddits:
        for term in search_terms[:3]:
            try:
                search_url = (
                    f"https://oauth.reddit.com/r/{sub}/search"
                    f"?q={requests.utils.quote(term)}"
                    f"&restrict_sr=1&sort=relevance&limit=20"
                )
                resp = requests.get(search_url, headers=headers, timeout=15)
                data = resp.json()
                children = data.get("data", {}).get("children", [])

                for post_wrapper in children:
                    p = post_wrapper.get("data", {})
                    text = (p.get("title", "") + " " +
                            p.get("selftext", "")).strip()
                    if len(text) > 60:
                        posts.append({
                            "text": text[:2000],
                            "username": p.get("author", "Reddit user"),
                            "platform": "reddit",
                            "source_url": f"https://reddit.com{p.get('permalink', '')}",
                            "date": datetime.fromtimestamp(
                                p.get("created_utc", 0)
                            ).isoformat() if p.get("created_utc") else "",
                        })
            except Exception:
                pass

    if progress_cb:
        progress_cb(f"Reddit: {len(posts)} posts collected")
    return posts


# ============================================================
# DATA FETCHER: Review Sites via Firecrawl
# ============================================================

def fetch_review_sites(firecrawl_key: str, brand: str, model: str,
                        progress_cb=None) -> list:
    posts = []
    brand_slug = brand.lower().replace(" ", "-")
    model_slug = model.lower().replace(" ", "-")

    review_urls = [
        # ZigWheels
        f"https://www.zigwheels.com/newcars/reviews/{brand_slug}-{model_slug}/user-reviews",
        # CarWale
        f"https://www.carwale.com/{brand_slug}-cars/{model_slug}/reviews/",
        # CarDekho
        f"https://www.cardekho.com/{brand_slug.replace('-', '')}-{model_slug}/user-reviews.htm",
        # Autocar India
        f"https://www.autocarindia.com/cars/{brand_slug}/{model_slug}/user-reviews",
        # GaadiWaadi
        f"https://www.gaadiwadi.com/{brand_slug}-{model_slug}/user-reviews",
        # Motorbeam
        f"https://www.motorbeam.com/?s={requests.utils.quote(model + ' problem')}",
        # V3Cars
        f"https://www.v3cars.com/{brand_slug}-{model_slug}/reviews",
        # MouthShut
        f"https://www.mouthshut.com/mobile-sites/search?keyword={requests.utils.quote(model)}",
        # CarTrade
        f"https://www.cartrade.com/{brand_slug}-cars/{model_slug}/user-reviews/",
        # Vicky.in
        f"https://www.vicky.in/{brand_slug}-{model_slug}/user-reviews",
        # Quora
        f"https://www.quora.com/search?q={requests.utils.quote(brand + ' ' + model + ' problems')}",
        # --- NEW SOURCES ADDED ---
        # Overdrive India — one of India's most established auto magazines
        f"https://www.overdriveindia.com/search?q={requests.utils.quote(brand + ' ' + model + ' problem')}",
        f"https://www.overdriveindia.com/cars/{brand_slug}/{model_slug}/user-reviews",
        # CarToq — practical ownership, maintenance, troubleshooting focus
        f"https://www.cartoq.com/{brand_slug}-{model_slug}/user-reviews",
        f"https://www.cartoq.com/search/?s={requests.utils.quote(model + ' problem')}",
        # AutoPortal — detailed user reviews
        f"https://www.autoportal.com/{brand_slug}-cars/{model_slug}/user-reviews/",
        # The Automotive India — active Indian community forum
        f"https://www.theautomotiveindia.com/forums/search/?q={requests.utils.quote(model)}&type=post",
        # BCMTouring — SUV and car ownership experience forum
        f"https://www.bcmtouring.com/forums/search/?q={requests.utils.quote(model + ' problem')}&type=post",
        # AutoX India — candid technical reviews
        f"https://www.autox.com/search/?q={requests.utils.quote(brand + ' ' + model + ' issues')}",
        # RSS feeds
        "https://www.autocarindia.com/rss/latest",
        "https://www.zigwheels.com/rss",
        "https://www.carwale.com/rss/news.xml",
        "https://www.overdriveindia.com/feed",
        "https://www.cartoq.com/feed",
        "https://www.motorbeam.com/feed",
    ]

    for i, url in enumerate(review_urls):
        try:
            resp = requests.post(
                "https://api.firecrawl.dev/v1/scrape",
                headers={
                    "Authorization": f"Bearer {firecrawl_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "url": url,
                    "formats": ["markdown"],
                    "onlyMainContent": True,
                },
                timeout=30,
            )
            data = resp.json()
            markdown = data.get("markdown", "")

            if not markdown:
                continue

            # For RSS feeds — filter articles mentioning model
            if "rss" in url or url.endswith(".xml"):
                if model.lower() not in markdown.lower():
                    continue

            blocks = re.split(r'\n---\n|\n\n\n|\n\n', markdown)
            for block in blocks:
                if len(block) > 80 and _has_quality_signal(block):
                    platform = (
                        "zigwheels" if "zigwheels" in url else
                        "carwale" if "carwale" in url else
                        "cardekho" if "cardekho" in url else
                        "autocar" if "autocar" in url else
                        "gaadiwadi" if "gaadiwadi" in url else
                        "motorbeam" if "motorbeam" in url else
                        "v3cars" if "v3cars" in url else
                        "mouthshut" if "mouthshut" in url else
                        "cartrade" if "cartrade" in url else
                        "vicky" if "vicky" in url else
                        "quora" if "quora" in url else
                        "overdrive" if "overdriveindia" in url else
                        "cartoq" if "cartoq" in url else
                        "autoportal" if "autoportal" in url else
                        "automotive_india" if "theautomotiveindia" in url else
                        "bcmtouring" if "bcmtouring" in url else
                        "autox" if "autox" in url else "review_site"
                    )
                    posts.append({
                        "text": block[:2000],
                        "username": _extract_username(block),
                        "platform": platform,
                        "source_url": url,
                        "date": data.get("metadata", {}).get(
                            "publishedTime", ""
                        ),
                    })

            if progress_cb and i % 3 == 0:
                progress_cb(
                    f"Review sites: {len(posts)} reviews collected ({i+1}/{len(review_urls)} sources)"
                )

        except Exception:
            pass

    if progress_cb:
        progress_cb(f"Review sites: {len(posts)} reviews total")
    return posts


# ============================================================
# FILTER: Quality defects only — exclude service/delivery
# ============================================================

EXCLUDE_PATTERNS = [
    r'service cent(re|er)', r'service advisor', r'waiting time',
    r'delivery delay', r'delayed delivery', r'waiting period',
    r'\bdealership\b', r'showroom staff', r'sales (person|executive|man)',
    r'test drive experience', r'booking cancel', r'booking amount',
    r'price hike', r'on-road price', r'ex-showroom', r'\binsurance\b',
    r'\bemi\b', r'\bloan\b', r'finance charge', r'discount offer',
    r'extended warranty cost',
]

INCLUDE_SIGNALS = [
    r'vibrat', r'rattle', r'\bnoise\b', r'squeak', r'creak', r'\bnvh\b',
    r'\bstall', r'\bjerk\b', r'hesitat', r'misfire', r'shudder',
    r'\bcrack\b', r'\bpeel\b', r'\brust\b', r'\bgap\b',
    r'fit (and|&) finish', r'\bleak', r'\bseep', r'\bsmell\b',
    r'\bsmoke\b', r'overheat', r'\bbrake', r'\bsteering', r'suspension',
    r'alignment', r'electrical', r'\bsensor', r'infotainment', r'\badas\b',
    r'\babs\b', r'\besp\b', r'\bairbag', r'\bcamera\b', r'\bproblem\b',
    r'\bissue\b', r'\bdefect\b', r'\bfault\b', r'\bfailure\b',
    r'\bbroke\b', r'broken', r'stopped working', r'not working',
    r'\bclunk\b', r'\bthud\b', r'\bgrind\b', r'\bscreech\b',
    r'hard to', r'difficult to', r'uncomfortable', r'blind spot',
    r'visibility', r'ergonomic', r'flickering', r'rattling', r'creaking',
    r'overheating', r'misfiring', r'stalling', r'pulling',
]

def hard_filter(comments: list) -> list:
    """Remove service/delivery noise. Keep only quality defects."""
    filtered = []
    for c in comments:
        text = c.get("text", "").lower()
        if len(text) < 40:
            continue
        # Exclude if matches any exclude pattern
        excluded = any(
            re.search(p, text) for p in EXCLUDE_PATTERNS
        )
        if excluded:
            continue
        # Include if matches any quality signal
        has_signal = any(
            re.search(p, text) for p in INCLUDE_SIGNALS
        )
        if has_signal:
            filtered.append(c)
    return filtered


def _has_quality_signal(text: str) -> bool:
    text_lower = text.lower()
    return any(re.search(p, text_lower) for p in INCLUDE_SIGNALS)


def _extract_username(text: str) -> str:
    patterns = [
        r'@([\w]+)',
        r'Posted by[:\s]+([\w\s]+)',
        r'Author[:\s]+([\w\s]+)',
        r'By[:\s]+([\w\s]{3,25})',
    ]
    for p in patterns:
        m = re.search(p, text, re.IGNORECASE)
        if m:
            return m.group(1).strip()
    return "Forum user"


# ============================================================
# DEDUPLICATION — Token-based Jaccard similarity
# ============================================================

def deduplicate(comments: list) -> list:
    """Cluster similar complaints together."""
    stop_words = {
        'the', 'is', 'at', 'which', 'on', 'and', 'or', 'but', 'in',
        'to', 'a', 'an', 'this', 'that', 'it', 'my', 'i', 'have',
        'has', 'been', 'was', 'are', 'for', 'of', 'with', 'by',
    }

    def get_words(text):
        return set(
            w for w in re.findall(r'\b\w+\b', text.lower())
            if len(w) > 4 and w not in stop_words
        )

    clusters = []
    for comment in comments:
        words = get_words(comment.get("text", ""))
        placed = False
        for cluster in clusters:
            rep_words = get_words(cluster[0].get("text", ""))
            intersection = len(words & rep_words)
            union = len(words | rep_words)
            similarity = intersection / union if union > 0 else 0
            if similarity > 0.28:
                cluster.append(comment)
                placed = True
                break
        if not placed:
            clusters.append([comment])

    result = []
    for cluster in clusters:
        platforms = list(set(c["platform"] for c in cluster))
        usernames = list(set(c.get("username", "") for c in cluster))
        result.append({
            **cluster[0],
            "mention_count": len(cluster),
            "all_mentions": cluster,
            "platforms": platforms,
            "usernames": usernames,
            "cross_platform_flag": len(platforms) >= 3,
            "cross_user_flag": len(usernames) >= 3,
        })

    return sorted(result, key=lambda x: x["mention_count"], reverse=True)


# ============================================================
# ESCALATION DETECTOR
# ============================================================

ESCALATION_KEYWORDS = [
    "consumer court", "consumer forum", "ncdrc", "cdrc",
    "district forum", "morth", "ministry of road",
    "replacement vehicle", "replace my car", "refund my car",
    "legal notice", "file a case", "advocate", "complained to company",
    "escalated to", "rto complaint", "lemon car",
    "demand replacement", "asking for replacement",
    "reported to government", "vahan complaint",
]

def detect_escalations(comments: list) -> list:
    """Find comments where customers escalated or demanded replacement."""
    escalations = []
    for c in comments:
        text = c.get("text", "").lower()
        matched = [kw for kw in ESCALATION_KEYWORDS if kw in text]
        if matched:
            esc_type = (
                "Consumer Forum" if any(
                    k in matched for k in [
                        "consumer court", "consumer forum",
                        "ncdrc", "cdrc", "district forum"
                    ]
                ) else
                "Govt Portal" if any(
                    k in matched for k in ["morth", "vahan", "rto"]
                ) else
                "Replacement Demand" if any(
                    k in matched for k in [
                        "replacement", "lemon", "refund"
                    ]
                ) else "Escalation"
            )
            escalations.append({
                "type": esc_type,
                "quote": c.get("text", "")[:500],
                "username": c.get("username", "Unknown"),
                "platform": c.get("platform", ""),
                "source_url": c.get("source_url", ""),
                "date": c.get("date", ""),
                "keywords": matched,
            })
    return escalations


# ============================================================
# GEMINI ANALYSIS — Zero hallucination
# ============================================================

GEMINI_SYSTEM_PROMPT = """You are AutoSentinel — a zero-hallucination automotive quality 
intelligence engine for an OEM R&D team.

YOUR PURPOSE: Extract real quality defects that customers experienced while owning and 
driving their vehicle.

ABSOLUTE RULES — NO EXCEPTIONS:
1. Extract ONLY information explicitly stated by customers. Never infer or elaborate.
2. If the condition under which a problem occurs is NOT stated: write exactly 
   "Condition not specified by customer" — never guess.
3. If mileage at onset is not mentioned: write "Mileage not specified"
4. EXCLUDE ENTIRELY: service wait times, delivery delays, dealership experience, 
   price issues, insurance, EMI, booking experience.
5. INCLUDE ONLY: mechanical defects, electrical faults, NVH (noise/vibration/harshness),
   build quality failures, fit and finish problems, ADAS malfunctions, handling defects,
   ergonomic difficulties experienced during ownership.
6. Every defect must be evidenced by actual quotes from the input. 
   Never create defects not evidenced in the text.
7. Rank defects by total mention count — highest first.
8. The quote field must be the customer's EXACT words — never paraphrase.
   "Customer reported vibration" is FORBIDDEN.
   "The steering shakes badly at 80kmph" is correct.
9. DVP recommendations must directly mirror customer-reported conditions only.
   Never add conditions customers did not mention.

RETURN ONLY VALID JSON. No markdown. No preamble. No explanation."""

def run_gemini_analysis(api_key: str, brand: str, model: str,
                         year: int, comments: list) -> Optional[dict]:
    """Send filtered comments to Gemini and get structured analysis."""

    BATCH_SIZE = 120
    all_defects = []
    last_proposal = None

    # Split into batches
    batches = [
        comments[i:i + BATCH_SIZE]
        for i in range(0, len(comments), BATCH_SIZE)
    ]
    if not batches:
        return None

    for batch_idx, batch in enumerate(batches):
        comments_text = "\n---\n".join([
            f"[{i+1}] Platform: {c.get('platform','?')} | "
            f"User: {c.get('username','?')} | "
            f"Date: {c.get('date','?')} | "
            f"URL: {c.get('source_url','?')}\n"
            f"Comment: {c.get('text','')}"
            for i, c in enumerate(batch)
        ])

        user_prompt = f"""Analyse these {len(batch)} customer comments about 
{brand} {model} {year}. Apply all rules strictly.

{comments_text}

Return JSON in this exact format:
{{
  "defects_found": [
    {{
      "title": "<max 8 words describing the defect>",
      "severity": "CRITICAL|HIGH|MODERATE",
      "mention_count": <integer>,
      "exact_condition": {{
        "when": "<exact situation stated by customer or null>",
        "terrain": "<road type if stated or null>",
        "onset_mileage": "<km range if stated or null>",
        "trigger": "<what makes it worse if stated or null>",
        "weather": "<climate context if stated or null>",
        "condition_completeness": "FULL|PARTIAL|NOT_SPECIFIED"
      }},
      "best_verbatim": {{
        "quote": "<exact customer words — verbatim>",
        "username": "<as posted>",
        "platform": "<source platform>",
        "source_url": "<url>",
        "date": "<if available>"
      }},
      "additional_verbatims": [
        {{"quote": "...", "username": "...", "platform": "...", "source_url": "...", "date": "..."}}
      ],
      "affected_component": "<specific part or system>",
      "sources": {{
        "team_bhp": 0, "youtube": 0, "reddit": 0, "zigwheels": 0,
        "carwale": 0, "cardekho": 0, "autocar": 0, "gaadiwadi": 0,
        "motorbeam": 0, "v3cars": 0, "quora": 0, "facebook": 0,
        "mouthshut": 0, "cartrade": 0, "vicky": 0,
        "overdrive": 0, "cartoq": 0, "autoportal": 0,
        "automotive_india": 0, "bcmtouring": 0, "autox": 0
      }},
      "cross_platform_flag": false,
      "cross_platform_count": 0
    }}
  ],
  "design_proposal": {{
    "management_summary": {{
      "customer_struggles": ["<plain language, no jargon>"],
      "next_model_changes": ["<plain language, no jargon>"]
    }},
    "technical_brief": [
      {{
        "defect_rank": 1,
        "root_cause_hypothesis": "<based strictly on customer reports only>",
        "affected_component": "<specific part>",
        "dvp_recommendation": "<what to test and exactly how>",
        "validation_condition": "<exact test condition mirroring customer reports>"
      }}
    ]
  }}
}}"""

        try:
            resp = requests.post(
                f"https://generativelanguage.googleapis.com/v1beta/models/"
                f"gemini-1.5-pro:generateContent?key={api_key}",
                headers={"Content-Type": "application/json"},
                json={
                    "system_instruction": {
                        "parts": [{"text": GEMINI_SYSTEM_PROMPT}]
                    },
                    "contents": [{"parts": [{"text": user_prompt}]}],
                    "generationConfig": {
                        "temperature": 0.1,
                        "maxOutputTokens": 8192,
                        "responseMimeType": "application/json",
                    },
                },
                timeout=90,
            )
            data = resp.json()
            raw_text = (
                data.get("candidates", [{}])[0]
                .get("content", {})
                .get("parts", [{}])[0]
                .get("text", "")
            )
            if raw_text:
                clean = raw_text.replace("```json", "").replace("```", "").strip()
                parsed = json.loads(clean)
                all_defects.extend(parsed.get("defects_found", []))
                if parsed.get("design_proposal"):
                    last_proposal = parsed["design_proposal"]

        except Exception as e:
            st.warning(f"Gemini batch {batch_idx + 1} error: {str(e)[:100]}")
            time.sleep(10)  # Rate limit backoff

    # Merge defects across batches
    merged = {}
    for d in all_defects:
        key = d["title"].lower().strip()
        if key in merged:
            merged[key]["mention_count"] += d.get("mention_count", 1)
            merged[key]["additional_verbatims"].extend(
                d.get("additional_verbatims", [])
            )
        else:
            merged[key] = d

    top10 = sorted(
        merged.values(),
        key=lambda x: x.get("mention_count", 0),
        reverse=True
    )[:10]

    for i, d in enumerate(top10):
        d["rank"] = i + 1

    return {
        "car": f"{brand} {model} {year}",
        "brand": brand,
        "model": model,
        "year": year,
        "total_comments_analysed": len(comments),
        "top_10_defects": top10,
        "design_proposal": last_proposal or {
            "management_summary": {
                "customer_struggles": [],
                "next_model_changes": [],
            },
            "technical_brief": [],
        },
    }


# ============================================================
# CACHE — Simple Streamlit session state cache
# ============================================================

def get_cache_key(brand: str, model: str, year: int) -> str:
    return hashlib.md5(
        f"{brand}_{model}_{year}".lower().encode()
    ).hexdigest()

def get_cached(cache_key: str) -> Optional[dict]:
    cache = st.session_state.get("analysis_cache", {})
    entry = cache.get(cache_key)
    if entry:
        cached_at = entry.get("cached_at")
        if cached_at:
            age = datetime.now() - datetime.fromisoformat(cached_at)
            if age < timedelta(hours=48):
                return entry.get("data")
    return None

def save_cache(cache_key: str, data: dict):
    if "analysis_cache" not in st.session_state:
        st.session_state["analysis_cache"] = {}
    st.session_state["analysis_cache"][cache_key] = {
        "data": data,
        "cached_at": datetime.now().isoformat(),
    }


# ============================================================
# EXPORT FUNCTIONS
# ============================================================

def export_pdf(result: dict, escalations: list) -> bytes:
    """Generate PDF report using reportlab."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, HRFlowable,
        )
        from reportlab.lib.units import mm
        import io

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer, pagesize=A4,
            rightMargin=20*mm, leftMargin=20*mm,
            topMargin=20*mm, bottomMargin=20*mm,
        )

        BG = colors.HexColor("#0A0A0A")
        RED = colors.HexColor("#E63946")
        TEAL = colors.HexColor("#2EC4B6")
        WHITE = colors.white
        GRAY = colors.HexColor("#A0A0A0")

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "Title", fontSize=22, textColor=WHITE,
            fontName="Helvetica-Bold", spaceAfter=8,
        )
        heading_style = ParagraphStyle(
            "Heading", fontSize=14, textColor=RED,
            fontName="Helvetica-Bold", spaceAfter=6,
        )
        body_style = ParagraphStyle(
            "Body", fontSize=10, textColor=WHITE,
            fontName="Helvetica", spaceAfter=4, leading=16,
        )
        small_style = ParagraphStyle(
            "Small", fontSize=8, textColor=GRAY,
            fontName="Helvetica", spaceAfter=3,
        )

        story = []

        # Cover
        story.append(Spacer(1, 30*mm))
        story.append(Paragraph("AutoSentinel", title_style))
        story.append(Paragraph(
            f"Quality Intelligence Brief — {result.get('car', '')}",
            heading_style
        ))
        story.append(Paragraph(
            f"Generated: {datetime.now().strftime('%d %B %Y, %H:%M')}",
            small_style
        ))
        story.append(Paragraph(
            "Confidential — Internal R&D Use Only", small_style
        ))
        story.append(Paragraph(
            "Service and delivery complaints excluded. Quality defects only.",
            small_style
        ))
        story.append(PageBreak())

        # Escalations
        if escalations:
            story.append(Paragraph("⚠ ESCALATION ALERTS", heading_style))
            story.append(Paragraph(
                f"{len(escalations)} issues escalated by customers",
                body_style
            ))
            for esc in escalations:
                story.append(Spacer(1, 4))
                story.append(Paragraph(
                    f"[{esc['type']}] {esc['username']} on {esc['platform']}",
                    small_style
                ))
                story.append(Paragraph(
                    f'"{esc["quote"][:300]}"', body_style
                ))
            story.append(PageBreak())

        # Top 10 defects
        for defect in result.get("top_10_defects", []):
            rank = defect.get("rank", "?")
            story.append(Paragraph(
                f"#{rank} — {defect.get('title', '')}",
                heading_style
            ))
            story.append(Paragraph(
                f"Severity: {defect.get('severity','?')} | "
                f"Mentions: {defect.get('mention_count', 0)}",
                small_style
            ))

            # Condition
            cond = defect.get("exact_condition", {})
            story.append(Paragraph("WHEN DOES IT OCCUR:", small_style))
            for field, label in [
                ("when", "When"), ("terrain", "Terrain"),
                ("onset_mileage", "Onset Mileage"),
                ("trigger", "Trigger"), ("weather", "Weather"),
            ]:
                val = cond.get(field) or "Not specified by customer"
                story.append(Paragraph(f"{label}: {val}", body_style))

            # Verbatim
            vb = defect.get("best_verbatim", {})
            if vb.get("quote"):
                story.append(Spacer(1, 4))
                story.append(Paragraph("CUSTOMER VERBATIM:", small_style))
                story.append(Paragraph(
                    f'"{vb["quote"]}"', body_style
                ))
                story.append(Paragraph(
                    f"— {vb.get('username','?')}, "
                    f"{vb.get('platform','?')}, "
                    f"{vb.get('date','?')}",
                    small_style
                ))

            story.append(HRFlowable(width="100%", thickness=0.5,
                                     color=GRAY))
            story.append(Spacer(1, 6))

        story.append(PageBreak())

        # Design proposal — management summary
        proposal = result.get("design_proposal", {})
        mgmt = proposal.get("management_summary", {})
        story.append(Paragraph("DESIGN PROPOSAL — Management Summary",
                               heading_style))
        story.append(Paragraph("What customers are struggling with:",
                               small_style))
        for item in mgmt.get("customer_struggles", []):
            story.append(Paragraph(f"• {item}", body_style))
        story.append(Spacer(1, 8))
        story.append(Paragraph("What must change in the next model:",
                               small_style))
        for item in mgmt.get("next_model_changes", []):
            story.append(Paragraph(f"• {item}", body_style))

        story.append(PageBreak())

        # Technical brief table
        tech = proposal.get("technical_brief", [])
        if tech:
            story.append(Paragraph("DESIGN PROPOSAL — Technical Brief",
                                   heading_style))
            table_data = [["Rank", "Component", "Root Cause Hypothesis",
                           "DVP Recommendation"]]
            for tb in tech:
                table_data.append([
                    str(tb.get("defect_rank", "")),
                    tb.get("affected_component", ""),
                    tb.get("root_cause_hypothesis", "")[:150],
                    tb.get("dvp_recommendation", "")[:150],
                ])
            t = Table(table_data, colWidths=[12*mm, 35*mm, 70*mm, 50*mm])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,0), RED),
                ("TEXTCOLOR", (0,0), (-1,0), WHITE),
                ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                ("FONTSIZE", (0,0), (-1,-1), 8),
                ("TEXTCOLOR", (0,1), (-1,-1), WHITE),
                ("BACKGROUND", (0,1), (-1,-1), BG),
                ("ROWBACKGROUNDS", (0,1), (-1,-1),
                 [colors.HexColor("#141414"), colors.HexColor("#1C1C1C")]),
                ("GRID", (0,0), (-1,-1), 0.3, GRAY),
                ("VALIGN", (0,0), (-1,-1), "TOP"),
                ("TOPPADDING", (0,0), (-1,-1), 4),
                ("BOTTOMPADDING", (0,0), (-1,-1), 4),
            ]))
            story.append(t)

        doc.build(story)
        return buffer.getvalue()

    except ImportError:
        st.warning("reportlab not installed. PDF export unavailable.")
        return b""


def export_excel(result: dict, escalations: list) -> bytes:
    """Generate Excel report using openpyxl."""
    try:
        import openpyxl
        from openpyxl.styles import (
            Font, PatternFill, Alignment, Border, Side,
        )
        import io

        wb = openpyxl.Workbook()

        # Colors
        RED_FILL = PatternFill("solid", fgColor="E63946")
        DARK_FILL = PatternFill("solid", fgColor="141414")
        CARD_FILL = PatternFill("solid", fgColor="1C1C1C")
        WHITE_FONT = Font(color="F0F0F0", bold=True)
        GRAY_FONT = Font(color="A0A0A0")
        RED_FONT = Font(color="E63946", bold=True)
        TEAL_FONT = Font(color="2EC4B6", bold=True)

        # Sheet 1: Summary
        ws1 = wb.active
        ws1.title = "Summary"
        ws1.sheet_view.showGridLines = False
        ws1["A1"] = "AutoSentinel — Quality Intelligence Brief"
        ws1["A1"].font = Font(color="E63946", bold=True, size=16)
        ws1["A2"] = result.get("car", "")
        ws1["A2"].font = Font(color="F0F0F0", bold=True, size=12)
        ws1["A3"] = f"Generated: {datetime.now().strftime('%d %B %Y %H:%M')}"
        ws1["A3"].font = GRAY_FONT
        ws1["A4"] = "Service and delivery complaints excluded. Quality defects only."
        ws1["A4"].font = Font(color="666666", italic=True)
        ws1["A6"] = "Total comments analysed:"
        ws1["B6"] = result.get("total_comments_analysed", 0)
        ws1["B6"].font = RED_FONT
        ws1["A7"] = "Total quality defects found:"
        ws1["B7"] = len(result.get("top_10_defects", []))
        ws1["A8"] = "Escalations detected:"
        ws1["B8"] = len(escalations)
        ws1["B8"].font = RED_FONT

        # Sheet 2: Top 10 Defects
        ws2 = wb.create_sheet("Top 10 Defects")
        ws2.sheet_view.showGridLines = False
        headers = [
            "Rank", "Severity", "Defect Title", "Mention Count",
            "Affected Component", "When", "Terrain", "Onset Mileage",
            "Trigger", "Weather", "Condition Completeness",
            "Best Verbatim Quote", "Username", "Platform", "Source URL",
            "Cross Platform Flag", "Sources Count",
        ]
        for col, h in enumerate(headers, 1):
            cell = ws2.cell(row=1, column=col, value=h)
            cell.fill = RED_FILL
            cell.font = WHITE_FONT
            cell.alignment = Alignment(wrap_text=True)

        for defect in result.get("top_10_defects", []):
            cond = defect.get("exact_condition", {})
            vb = defect.get("best_verbatim", {})
            sources = defect.get("sources", {})
            row_data = [
                defect.get("rank", ""),
                defect.get("severity", ""),
                defect.get("title", ""),
                defect.get("mention_count", 0),
                defect.get("affected_component", ""),
                cond.get("when") or "Not specified by customer",
                cond.get("terrain") or "Not specified",
                cond.get("onset_mileage") or "Not specified",
                cond.get("trigger") or "Not specified",
                cond.get("weather") or "Not specified",
                cond.get("condition_completeness", "NOT_SPECIFIED"),
                vb.get("quote", ""),
                vb.get("username", ""),
                vb.get("platform", ""),
                vb.get("source_url", ""),
                "Yes" if defect.get("cross_platform_flag") else "No",
                sum(sources.values()),
            ]
            row_num = defect.get("rank", 1) + 1
            for col, val in enumerate(row_data, 1):
                cell = ws2.cell(row=row_num, column=col, value=val)
                cell.font = Font(color="F0F0F0")
                cell.fill = CARD_FILL if row_num % 2 == 0 else DARK_FILL
                cell.alignment = Alignment(wrap_text=True, vertical="top")

        # Sheet 3: All Verbatims
        ws3 = wb.create_sheet("All Verbatims")
        ws3.sheet_view.showGridLines = False
        v_headers = ["Defect Rank", "Defect Title", "Quote",
                     "Username", "Platform", "Source URL", "Date"]
        for col, h in enumerate(v_headers, 1):
            cell = ws3.cell(row=1, column=col, value=h)
            cell.fill = RED_FILL
            cell.font = WHITE_FONT

        v_row = 2
        for defect in result.get("top_10_defects", []):
            vb = defect.get("best_verbatim", {})
            ws3.cell(row=v_row, column=1,
                     value=defect.get("rank", "")).font = RED_FONT
            ws3.cell(row=v_row, column=2,
                     value=defect.get("title", "")).font = Font(
                color="F0F0F0"
            )
            ws3.cell(row=v_row, column=3,
                     value=vb.get("quote", "")).font = Font(
                color="E8E8E8", italic=True
            )
            ws3.cell(row=v_row, column=4,
                     value=vb.get("username", "")).font = GRAY_FONT
            ws3.cell(row=v_row, column=5,
                     value=vb.get("platform", "")).font = GRAY_FONT
            ws3.cell(row=v_row, column=6,
                     value=vb.get("source_url", "")).font = GRAY_FONT
            ws3.cell(row=v_row, column=7,
                     value=vb.get("date", "")).font = GRAY_FONT
            v_row += 1
            for av in defect.get("additional_verbatims", []):
                ws3.cell(row=v_row, column=1,
                         value=defect.get("rank", "")).font = GRAY_FONT
                ws3.cell(row=v_row, column=2, value="").font = GRAY_FONT
                ws3.cell(row=v_row, column=3,
                         value=av.get("quote", "")).font = Font(
                    color="999999", italic=True
                )
                ws3.cell(row=v_row, column=4,
                         value=av.get("username", "")).font = GRAY_FONT
                ws3.cell(row=v_row, column=5,
                         value=av.get("platform", "")).font = GRAY_FONT
                ws3.cell(row=v_row, column=6,
                         value=av.get("source_url", "")).font = GRAY_FONT
                ws3.cell(row=v_row, column=7,
                         value=av.get("date", "")).font = GRAY_FONT
                v_row += 1

        # Sheet 4: Escalations
        ws4 = wb.create_sheet("Escalations")
        ws4.sheet_view.showGridLines = False
        esc_headers = ["Type", "Quote", "Username",
                       "Platform", "Source URL", "Date", "Keywords"]
        for col, h in enumerate(esc_headers, 1):
            cell = ws4.cell(row=1, column=col, value=h)
            cell.fill = RED_FILL
            cell.font = WHITE_FONT

        for i, esc in enumerate(escalations, 2):
            ws4.cell(row=i, column=1, value=esc.get("type", "")).font = RED_FONT
            ws4.cell(row=i, column=2, value=esc.get("quote", "")).font = Font(color="F0F0F0")
            ws4.cell(row=i, column=3, value=esc.get("username", "")).font = GRAY_FONT
            ws4.cell(row=i, column=4, value=esc.get("platform", "")).font = GRAY_FONT
            ws4.cell(row=i, column=5, value=esc.get("source_url", "")).font = GRAY_FONT
            ws4.cell(row=i, column=6, value=esc.get("date", "")).font = GRAY_FONT
            ws4.cell(row=i, column=7, value=", ".join(esc.get("keywords", []))).font = GRAY_FONT

        # Sheet 5: Design Proposal
        ws5 = wb.create_sheet("Design Proposal")
        ws5.sheet_view.showGridLines = False
        proposal = result.get("design_proposal", {})
        mgmt = proposal.get("management_summary", {})
        ws5["A1"] = "MANAGEMENT SUMMARY"
        ws5["A1"].font = Font(color="2EC4B6", bold=True, size=13)
        ws5["A2"] = "What customers are struggling with:"
        ws5["A2"].font = Font(color="A0A0A0", italic=True)
        for i, item in enumerate(mgmt.get("customer_struggles", []), 3):
            ws5.cell(row=i, column=1, value=f"• {item}").font = Font(color="F0F0F0")
        offset = len(mgmt.get("customer_struggles", [])) + 4
        ws5.cell(row=offset, column=1,
                 value="What must change in the next model:").font = Font(
            color="A0A0A0", italic=True
        )
        for i, item in enumerate(mgmt.get("next_model_changes", []), offset + 1):
            ws5.cell(row=i, column=1, value=f"• {item}").font = Font(color="F0F0F0")

        tech_start = offset + len(mgmt.get("next_model_changes", [])) + 3
        ws5.cell(row=tech_start, column=1,
                 value="TECHNICAL BRIEF").font = Font(
            color="2EC4B6", bold=True, size=13
        )
        tech_headers = ["Rank", "Component", "Root Cause Hypothesis",
                        "DVP Recommendation", "Validation Condition"]
        for col, h in enumerate(tech_headers, 1):
            cell = ws5.cell(row=tech_start + 1, column=col, value=h)
            cell.fill = RED_FILL
            cell.font = WHITE_FONT

        for i, tb in enumerate(proposal.get("technical_brief", []),
                                tech_start + 2):
            ws5.cell(row=i, column=1, value=tb.get("defect_rank", "")).font = RED_FONT
            ws5.cell(row=i, column=2, value=tb.get("affected_component", "")).font = Font(color="F0F0F0")
            ws5.cell(row=i, column=3, value=tb.get("root_cause_hypothesis", "")).font = Font(color="F0F0F0")
            ws5.cell(row=i, column=4, value=tb.get("dvp_recommendation", "")).font = Font(color="F0F0F0")
            ws5.cell(row=i, column=5, value=tb.get("validation_condition", "")).font = Font(color="F0F0F0")

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    except ImportError:
        st.warning("openpyxl not installed. Excel export unavailable.")
        return b""


def export_pptx(result: dict, escalations: list) -> bytes:
    """Generate PowerPoint using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        BG_COLOR = RGBColor(0x0A, 0x0A, 0x0A)
        RED_COLOR = RGBColor(0xE6, 0x39, 0x46)
        TEAL_COLOR = RGBColor(0x2E, 0xC4, 0xB6)
        WHITE_COLOR = RGBColor(0xF0, 0xF0, 0xF0)
        GRAY_COLOR = RGBColor(0x66, 0x66, 0x66)

        blank_layout = prs.slide_layouts[6]

        def set_bg(slide):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR

        def add_text(slide, text, left, top, width, height,
                     size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                     italic=False):
            tb = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color or WHITE_COLOR
            return tb

        # Slide 1: Title
        slide1 = prs.slides.add_slide(blank_layout)
        set_bg(slide1)
        add_text(slide1, "AUTOSENTINEL", 0.5, 1.5, 12, 1.2,
                 size=36, bold=True, color=RED_COLOR)
        add_text(slide1, f"Quality Intelligence Brief — {result.get('car','')}",
                 0.5, 2.8, 12, 0.8, size=20, color=WHITE_COLOR)
        add_text(slide1, f"Generated: {datetime.now().strftime('%d %B %Y')}",
                 0.5, 3.7, 6, 0.5, size=12, color=GRAY_COLOR)
        add_text(slide1, "Confidential — Internal R&D Use Only",
                 0.5, 4.1, 8, 0.5, size=12, color=GRAY_COLOR, italic=True)
        add_text(slide1, "Service and delivery complaints excluded. Quality defects only.",
                 0.5, 4.5, 10, 0.5, size=11, color=GRAY_COLOR, italic=True)

        # Slide 2: Summary stats
        slide2 = prs.slides.add_slide(blank_layout)
        set_bg(slide2)
        add_text(slide2, "Executive Summary", 0.5, 0.3, 12, 0.7,
                 size=24, bold=True, color=RED_COLOR)
        add_text(slide2, f"{result.get('total_comments_analysed', 0):,}",
                 0.5, 1.2, 4, 1.2, size=48, bold=True, color=RED_COLOR)
        add_text(slide2, "Quality comments analysed",
                 0.5, 2.3, 4, 0.5, size=12, color=GRAY_COLOR)
        add_text(slide2, f"{len(result.get('top_10_defects', []))}",
                 5, 1.2, 3, 1.2, size=48, bold=True, color=TEAL_COLOR)
        add_text(slide2, "Defects identified",
                 5, 2.3, 3, 0.5, size=12, color=GRAY_COLOR)
        add_text(slide2, f"{len(escalations)}",
                 9, 1.2, 3, 1.2, size=48, bold=True,
                 color=RED_COLOR if escalations else GRAY_COLOR)
        add_text(slide2, "Customer escalations",
                 9, 2.3, 3, 0.5, size=12, color=GRAY_COLOR)

        # Top 3 issues preview
        add_text(slide2, "Top 3 Issues:", 0.5, 3.2, 12, 0.5,
                 size=13, bold=True, color=WHITE_COLOR)
        for i, defect in enumerate(result.get("top_10_defects", [])[:3]):
            add_text(slide2,
                     f"#{defect.get('rank','?')} {defect.get('title','')} "
                     f"— {defect.get('mention_count',0)} mentions",
                     0.5, 3.7 + i * 0.5, 12, 0.5, size=12, color=WHITE_COLOR)

        # Slides 3-12: One per defect
        for defect in result.get("top_10_defects", []):
            slide = prs.slides.add_slide(blank_layout)
            set_bg(slide)

            sev = defect.get("severity", "MODERATE")
            sev_color = (
                RED_COLOR if sev == "CRITICAL" else
                RGBColor(0xF4, 0xA2, 0x61) if sev == "HIGH" else
                RGBColor(0x52, 0xB7, 0x88)
            )

            add_text(slide,
                     f"#{defect.get('rank','?')} of 10   [{sev}]   "
                     f"{defect.get('mention_count',0)} mentions",
                     0.3, 0.2, 12, 0.4, size=10, color=GRAY_COLOR)
            add_text(slide, defect.get("title", ""), 0.3, 0.6, 12, 0.8,
                     size=22, bold=True, color=WHITE_COLOR)

            # Condition box
            cond = defect.get("exact_condition", {})
            add_text(slide, "WHEN DOES IT OCCUR",
                     0.3, 1.5, 6, 0.3, size=9, bold=True, color=RED_COLOR)
            cond_lines = []
            for field, label in [
                ("when", "When"), ("terrain", "Terrain"),
                ("onset_mileage", "Mileage"), ("trigger", "Trigger"),
            ]:
                val = cond.get(field) or "Not specified by customer"
                cond_lines.append(f"{label}: {val}")
            add_text(slide, "\n".join(cond_lines), 0.3, 1.8, 6, 1.8,
                     size=11, color=WHITE_COLOR)

            # Verbatim
            vb = defect.get("best_verbatim", {})
            add_text(slide, "CUSTOMER VERBATIM",
                     6.8, 1.5, 6, 0.3, size=9, bold=True, color=GRAY_COLOR)
            quote = vb.get("quote", "")[:280]
            add_text(slide, f'"{quote}"', 6.8, 1.8, 6, 1.5,
                     size=11, italic=True, color=WHITE_COLOR)
            add_text(slide,
                     f"— {vb.get('username','?')}, "
                     f"{vb.get('platform','?')}, {vb.get('date','?')}",
                     6.8, 3.3, 6, 0.4, size=9, color=GRAY_COLOR)

            # Component
            add_text(slide,
                     f"Affected: {defect.get('affected_component','')}",
                     0.3, 3.6, 12, 0.4, size=11, color=TEAL_COLOR)

        # Slide: Escalations (if any)
        if escalations:
            slide_esc = prs.slides.add_slide(blank_layout)
            set_bg(slide_esc)
            add_text(slide_esc, "⚠ ESCALATION ALERTS",
                     0.3, 0.2, 12, 0.6, size=22, bold=True, color=RED_COLOR)
            add_text(slide_esc,
                     f"{len(escalations)} issues escalated by customers",
                     0.3, 0.9, 12, 0.4, size=13, color=WHITE_COLOR)
            for i, esc in enumerate(escalations[:4]):
                add_text(slide_esc,
                         f"[{esc['type']}] {esc['username']} — "
                         f"{esc['quote'][:180]}",
                         0.3, 1.5 + i * 1.3, 12, 1.1,
                         size=11, color=WHITE_COLOR)

        # Slide: Management summary
        slide_mgmt = prs.slides.add_slide(blank_layout)
        set_bg(slide_mgmt)
        add_text(slide_mgmt, "Design Proposal — Management Summary",
                 0.3, 0.2, 12, 0.6, size=20, bold=True, color=TEAL_COLOR)
        proposal = result.get("design_proposal", {})
        mgmt = proposal.get("management_summary", {})
        add_text(slide_mgmt, "What customers are struggling with:",
                 0.3, 1.0, 12, 0.4, size=12, color=GRAY_COLOR)
        struggles_text = "\n".join([
            f"• {s}" for s in mgmt.get("customer_struggles", [])
        ])
        add_text(slide_mgmt, struggles_text, 0.3, 1.4, 6, 3,
                 size=11, color=WHITE_COLOR)
        add_text(slide_mgmt, "What must change in the next model:",
                 6.8, 1.0, 6, 0.4, size=12, color=GRAY_COLOR)
        changes_text = "\n".join([
            f"• {s}" for s in mgmt.get("next_model_changes", [])
        ])
        add_text(slide_mgmt, changes_text, 6.8, 1.4, 6, 3,
                 size=11, color=WHITE_COLOR)

        # Slide: Technical brief
        slide_tech = prs.slides.add_slide(blank_layout)
        set_bg(slide_tech)
        add_text(slide_tech, "Design Proposal — Technical Brief",
                 0.3, 0.2, 12, 0.6, size=20, bold=True, color=TEAL_COLOR)
        for i, tb in enumerate(proposal.get("technical_brief", [])[:5]):
            add_text(slide_tech,
                     f"#{tb.get('defect_rank','')} {tb.get('affected_component','')}:"
                     f" {tb.get('dvp_recommendation','')[:120]}",
                     0.3, 0.9 + i * 1.1, 12, 1.0,
                     size=11, color=WHITE_COLOR)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    except ImportError:
        st.warning("python-pptx not installed. PPTX export unavailable.")
        return b""


# ============================================================
# UI COMPONENTS
# ============================================================

def render_header():
    st.markdown("""
    <div class="as-header">
      <div>
        <span class="as-logo">
          <span class="red">AUTO</span><span class="white">SENTINEL</span>
        </span><br/>
        <span class="as-tagline">Customer Voice Intelligence · Quality Defects Only</span>
      </div>
      <div style="text-align:right">
        <span class="as-tagline">Internal R&amp;D Platform · Powered by Gemini AI</span>
      </div>
    </div>
    """, unsafe_allow_html=True)


def render_severity_badge(severity: str) -> str:
    cls = (
        "severity-critical" if severity == "CRITICAL" else
        "severity-high" if severity == "HIGH" else
        "severity-moderate"
    )
    return f'<span class="{cls}">{severity}</span>'


def render_source_badges(sources: dict) -> str:
    html = ""
    for source, count in sources.items():
        if count and count > 0:
            html += (
                f'<span class="source-badge">'
                f'{source.replace("_"," ").title()} {count}'
                f'</span>'
            )
    return html


def render_condition_box(cond: dict) -> str:
    def val(field, label):
        v = cond.get(field)
        val_class = "condition-val" if v else "condition-null"
        display = v if v else "Not specified by customer"
        return (
            f'<div class="condition-row">'
            f'<span class="condition-key">{label}</span>'
            f'<span class="{val_class}">{display}</span>'
            f'</div>'
        )

    completeness = cond.get("condition_completeness", "NOT_SPECIFIED")
    comp_color = (
        "#52B788" if completeness == "FULL" else
        "#F4A261" if completeness == "PARTIAL" else "#555"
    )

    return f"""
    <div class="condition-box">
      <div class="condition-label">When does it occur</div>
      {val("when", "When")}
      {val("terrain", "Terrain")}
      {val("onset_mileage", "Onset mileage")}
      {val("trigger", "Trigger")}
      {val("weather", "Weather / climate")}
      <div style="margin-top:8px">
        <span style="font-size:10px;color:{comp_color};
          font-family:'IBM Plex Mono',monospace;letter-spacing:1px;">
          CONDITION: {completeness}
        </span>
      </div>
    </div>
    """


def render_verbatim(vb: dict, label="CUSTOMER VERBATIM") -> str:
    quote = vb.get("quote", "")
    username = vb.get("username", "")
    platform = vb.get("platform", "")
    date = vb.get("date", "")
    url = vb.get("source_url", "")
    link = (
        f'<a href="{url}" target="_blank" '
        f'style="color:#444;font-size:10px;">[source]</a>'
        if url else ""
    )
    return f"""
    <div class="verbatim-box">
      <div class="verbatim-label">{label}</div>
      <div class="verbatim-quote">"{quote}"</div>
      <div class="verbatim-attr">
        — {username} &nbsp;·&nbsp; {platform} &nbsp;·&nbsp; {date}
        &nbsp;{link}
      </div>
    </div>
    """


def render_defect_card(defect: dict, idx: int):
    rank = defect.get("rank", idx + 1)
    title = defect.get("title", "")
    severity = defect.get("severity", "MODERATE")
    mention_count = defect.get("mention_count", 0)
    cond = defect.get("exact_condition", {})
    vb = defect.get("best_verbatim", {})
    sources = defect.get("sources", {})
    cross = defect.get("cross_platform_flag", False)
    cross_count = defect.get("cross_platform_count", 0)
    additional_vbs = defect.get("additional_verbatims", [])
    component = defect.get("affected_component", "")

    severity_html = render_severity_badge(severity)
    condition_html = render_condition_box(cond)
    verbatim_html = render_verbatim(vb)
    sources_html = render_source_badges(sources)

    cross_flag_html = ""
    if cross:
        cross_flag_html = (
            f'<div style="margin-top:8px;font-size:12px;color:#F4A261;">'
            f'🔁 Same defect reported across multiple platforms'
            f'</div>'
        )

    st.markdown(f"""
    <div class="defect-card">
      <div style="display:flex;justify-content:space-between;align-items:center">
        <span class="defect-rank">#{rank} of 10</span>
        <span>
          {severity_html}
          &nbsp;&nbsp;
          <span class="mention-count">{mention_count} mentions</span>
        </span>
      </div>
      <div class="defect-title">{title}</div>
      {f'<div style="font-size:12px;color:#2EC4B6;margin-bottom:12px;">Affected: {component}</div>' if component else ''}
      {condition_html}
      {verbatim_html}
      <div style="margin-top:10px">{sources_html}</div>
      {cross_flag_html}
    </div>
    """, unsafe_allow_html=True)

    # Additional verbatims in expander
    if additional_vbs:
        with st.expander(f"View {len(additional_vbs)} more customer quotes"):
            for av in additional_vbs[:10]:
                st.markdown(
                    render_verbatim(av, "ADDITIONAL VERBATIM"),
                    unsafe_allow_html=True
                )


def render_escalations(escalations: list):
    if not escalations:
        return
    st.markdown(f"""
    <div class="escalation-banner">
      <div class="escalation-title">
        ⚠ ESCALATION ALERT — {len(escalations)} Critical Issues
      </div>
      <div style="font-size:13px;color:#A0A0A0;margin-bottom:12px;">
        These customers have escalated beyond normal feedback channels.
        Requires immediate attention.
      </div>
    """, unsafe_allow_html=True)

    for esc in escalations:
        st.markdown(f"""
      <div class="escalation-card">
        <span class="escalation-type">{esc.get('type','ESCALATION')}</span>
        <div class="verbatim-quote" style="font-size:13px;">
          "{esc.get('quote','')[:400]}"
        </div>
        <div class="verbatim-attr">
          {esc.get('username','')} · {esc.get('platform','')} · 
          {esc.get('date','')}
          {'<a href="' + esc.get('source_url','') + '" target="_blank" style="color:#444;font-size:10px;margin-left:6px;">[source]</a>' if esc.get('source_url') else ''}
        </div>
        <div style="font-size:11px;color:#E63946;margin-top:6px;">
          Keywords: {', '.join(esc.get('keywords',[]))}
        </div>
      </div>
        """, unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)


def render_design_proposal(proposal: dict):
    st.markdown(
        '<div class="as-section-title">Design Proposal for Next Model</div>',
        unsafe_allow_html=True
    )

    tab1, tab2 = st.tabs(["📋 Management Summary", "🔧 Technical Brief"])

    with tab1:
        mgmt = proposal.get("management_summary", {})
        st.markdown('<div class="proposal-box">', unsafe_allow_html=True)
        st.markdown(
            '<div class="proposal-title">What customers are struggling with</div>',
            unsafe_allow_html=True
        )
        for item in mgmt.get("customer_struggles", []):
            st.markdown(f"• {item}")

        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown(
            '<div class="proposal-title">What must change in the next model</div>',
            unsafe_allow_html=True
        )
        for item in mgmt.get("next_model_changes", []):
            st.markdown(f"• {item}")
        st.markdown('</div>', unsafe_allow_html=True)

    with tab2:
        tech = proposal.get("technical_brief", [])
        if tech:
            import pandas as pd
            rows = []
            for tb in tech:
                rows.append({
                    "Rank": tb.get("defect_rank", ""),
                    "Component": tb.get("affected_component", ""),
                    "Root Cause Hypothesis": tb.get(
                        "root_cause_hypothesis", ""
                    ),
                    "DVP Recommendation": tb.get(
                        "dvp_recommendation", ""
                    ),
                    "Validation Condition": tb.get(
                        "validation_condition", ""
                    ),
                })
            df = pd.DataFrame(rows)
            st.dataframe(df, use_container_width=True, hide_index=True)
        else:
            st.info("Technical brief not available for this analysis.")


def render_summary_strip(result: dict, escalations: list,
                          cached: bool = False):
    total = result.get("total_comments_analysed", 0)
    car = result.get("car", "")
    freshness = (
        "🟢 Live · Just fetched"
        if not cached
        else "🟡 Cached · < 48 hours old"
    )

    col1, col2, col3, col4 = st.columns([3, 2, 2, 2])
    with col1:
        st.markdown(
            f'<div class="car-name-display">{car}</div>',
            unsafe_allow_html=True
        )
        st.markdown(
            '<span class="disclaimer-pill">'
            'Service &amp; delivery excluded · Quality defects only'
            '</span>',
            unsafe_allow_html=True
        )
    with col2:
        st.markdown(
            f'<div class="big-number">{total:,}</div>'
            f'<div class="big-number-label">Comments analysed</div>',
            unsafe_allow_html=True
        )
    with col3:
        st.markdown(
            f'<div class="big-number" style="color:#2EC4B6">'
            f'{len(result.get("top_10_defects", []))}</div>'
            f'<div class="big-number-label">Defects identified</div>',
            unsafe_allow_html=True
        )
    with col4:
        esc_color = "#E63946" if escalations else "#444"
        st.markdown(
            f'<div class="big-number" style="color:{esc_color}">'
            f'{len(escalations)}</div>'
            f'<div class="big-number-label">Escalations detected</div>',
            unsafe_allow_html=True
        )

    st.markdown(
        f'<div style="font-size:11px;color:#444;margin-top:8px;'
        f'font-family:IBM Plex Mono,monospace;">{freshness}</div>',
        unsafe_allow_html=True
    )
    st.markdown("<hr style='border-color:#1A1A1A;margin:16px 0'>",
                unsafe_allow_html=True)


# ============================================================
# MAIN PIPELINE
# ============================================================

def run_analysis(brand: str, model: str, year: int, secrets: dict,
                  force_refresh: bool = False):
    """Full pipeline: fetch → filter → dedup → escalation → gemini."""

    cache_key = get_cache_key(brand, model, year)

    if not force_refresh:
        cached = get_cached(cache_key)
        if cached:
            return cached, True  # (result, is_cached)

    progress_placeholder = st.empty()
    log_lines = []

    def log(msg):
        log_lines.append(f"✓ {msg}")
        progress_placeholder.markdown(
            "\n\n".join(log_lines[-8:])
        )

    all_comments = []

    # Step 1: YouTube
    with st.spinner(""):
        st.markdown(
            '<div style="color:#E63946;font-family:IBM Plex Mono,'
            'monospace;font-size:12px;">Step 1/10 — YouTube...</div>',
            unsafe_allow_html=True
        )
        yt = fetch_youtube(secrets["youtube"], brand, model, year,
                            progress_cb=log)
        all_comments.extend(yt)
        log(f"YouTube: {len(yt)} comments")

    # Step 2: Team-BHP
    st.markdown(
        '<div style="color:#E63946;font-family:IBM Plex Mono,'
        'monospace;font-size:12px;">Step 2/10 — Team-BHP...</div>',
        unsafe_allow_html=True
    )
    tbhp = fetch_teambhp(secrets["firecrawl"], brand, model,
                          progress_cb=log)
    all_comments.extend(tbhp)
    log(f"Team-BHP: {len(tbhp)} posts")

    # Step 3: Reddit
    st.markdown(
        '<div style="color:#E63946;font-family:IBM Plex Mono,'
        'monospace;font-size:12px;">Step 3/10 — Reddit...</div>',
        unsafe_allow_html=True
    )
    reddit = fetch_reddit(
        secrets["reddit_id"], secrets["reddit_secret"],
        brand, model, progress_cb=log
    )
    all_comments.extend(reddit)
    log(f"Reddit: {len(reddit)} posts")

    # Steps 4-7: Review sites
    st.markdown(
        '<div style="color:#E63946;font-family:IBM Plex Mono,'
        'monospace;font-size:12px;">'
        'Steps 4-7/10 — Review sites...</div>',
        unsafe_allow_html=True
    )
    reviews = fetch_review_sites(secrets["firecrawl"], brand, model,
                                  progress_cb=log)
    all_comments.extend(reviews)
    log(f"Review sites: {len(reviews)} reviews")

    # Step 8: Hard filter
    st.markdown(
        '<div style="color:#E63946;font-family:IBM Plex Mono,'
        'monospace;font-size:12px;">Step 8/10 — Filtering...</div>',
        unsafe_allow_html=True
    )
    filtered = hard_filter(all_comments)
    log(f"Filter: {len(filtered)} quality defect comments kept "
        f"from {len(all_comments)} total")

    if len(filtered) < 20:
        st.markdown(
            '<div class="warning-box">⚠ Limited quality data found '
            f'({len(filtered)} comments). Results may not be fully '
            f'representative. Try a more popular model or year.</div>',
            unsafe_allow_html=True
        )

    # Step 9: Deduplication + Escalation detection
    st.markdown(
        '<div style="color:#E63946;font-family:IBM Plex Mono,'
        'monospace;font-size:12px;">'
        'Step 9/10 — Deduplication &amp; escalation scan...</div>',
        unsafe_allow_html=True
    )
    deduped = deduplicate(filtered)
    escalations = detect_escalations(all_comments)
    log(f"Deduplication: {len(deduped)} unique defect clusters")
    log(f"Escalations detected: {len(escalations)}")

    # Step 10: Gemini
    st.markdown(
        '<div style="color:#E63946;font-family:IBM Plex Mono,'
        'monospace;font-size:12px;">'
        'Step 10/10 — Gemini AI analysis...</div>',
        unsafe_allow_html=True
    )
    result = run_gemini_analysis(
        secrets["gemini"], brand, model, year, deduped
    )

    if not result:
        st.error("Analysis failed. Please try again.")
        return None, False

    log(f"Analysis complete. {len(result.get('top_10_defects',[]))} "
        f"defects identified.")

    # Store escalations in result for export
    result["escalations"] = escalations

    progress_placeholder.empty()
    save_cache(cache_key, {"result": result, "escalations": escalations})
    return {"result": result, "escalations": escalations}, False


# ============================================================
# MAIN APP
# ============================================================

def main():
    render_header()

    secrets = get_secrets()

    # Mode toggle
    mode = st.radio(
        "Analysis mode",
        ["🔍 Single Model Analysis", "⚖️ Compare Two Models"],
        horizontal=True,
        label_visibility="collapsed",
    )
    st.markdown("<br>", unsafe_allow_html=True)

    # ---- SINGLE MODE ----
    if "Single" in mode:
        col1, col2, col3 = st.columns([2, 2, 1])
        with col1:
            brand = st.selectbox(
                "Brand",
                [""] + list(INDIAN_CAR_MODELS.keys()),
                key="brand_single"
            )
        with col2:
            model_options = (
                [""] + INDIAN_CAR_MODELS.get(brand, [])
                if brand else [""]
            )
            model = st.selectbox("Model", model_options, key="model_single")
        with col3:
            year = st.selectbox(
                "Year",
                list(range(2025, 2017, -1)),
                key="year_single"
            )

        manual = st.text_input(
            "Or type any model manually (overrides dropdown)",
            placeholder="e.g. Kia Carens 2023",
            key="manual_single"
        )

        if manual.strip():
            parts = manual.strip().split()
            if len(parts) >= 2:
                brand = parts[0]
                model = " ".join(parts[1:])

        col_btn, col_refresh = st.columns([3, 1])
        with col_btn:
            analyse_btn = st.button(
                "🔍 Analyse Customer Voice",
                use_container_width=True,
                key="analyse_single"
            )
        with col_refresh:
            force_refresh = st.checkbox("Force refresh", key="refresh_single")

        if analyse_btn and brand and model:
            cache_key = get_cache_key(brand, model, year)
            cached_data = None if force_refresh else get_cached(cache_key)

            if cached_data:
                st.success("⚡ Loaded from cache (< 48 hrs old)")
                data = cached_data
                is_cached = True
            else:
                with st.status(
                    f"Analysing {brand} {model} {year}...",
                    expanded=True
                ) as status:
                    data, is_cached = run_analysis(
                        brand, model, year, secrets,
                        force_refresh=force_refresh
                    )
                    if data:
                        status.update(
                            label="✅ Analysis complete!", state="complete"
                        )

            if data:
                result = data["result"]
                escalations = data["escalations"]

                render_summary_strip(result, escalations, is_cached)
                render_escalations(escalations)

                st.markdown(
                    '<div class="as-section-title">Top 10 Quality Defects</div>',
                    unsafe_allow_html=True
                )

                for i, defect in enumerate(
                    result.get("top_10_defects", [])
                ):
                    render_defect_card(defect, i)

                render_design_proposal(result.get("design_proposal", {}))

                # Exports
                st.markdown(
                    '<div class="as-section-title">Export Report</div>',
                    unsafe_allow_html=True
                )
                exp_col1, exp_col2, exp_col3 = st.columns(3)

                with exp_col1:
                    if st.button("📄 Download PDF", use_container_width=True):
                        pdf_bytes = export_pdf(result, escalations)
                        if pdf_bytes:
                            st.download_button(
                                "⬇ Save PDF",
                                data=pdf_bytes,
                                file_name=f"autosentinel_{brand}_{model}_{year}.pdf",
                                mime="application/pdf",
                                use_container_width=True,
                            )

                with exp_col2:
                    if st.button("📊 Download Excel", use_container_width=True):
                        xl_bytes = export_excel(result, escalations)
                        if xl_bytes:
                            st.download_button(
                                "⬇ Save Excel",
                                data=xl_bytes,
                                file_name=f"autosentinel_{brand}_{model}_{year}.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                use_container_width=True,
                            )

                with exp_col3:
                    if st.button("📑 Download PPTX", use_container_width=True):
                        pptx_bytes = export_pptx(result, escalations)
                        if pptx_bytes:
                            st.download_button(
                                "⬇ Save PowerPoint",
                                data=pptx_bytes,
                                file_name=f"autosentinel_{brand}_{model}_{year}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True,
                            )

        elif analyse_btn:
            st.warning("Please select a brand and model to analyse.")

    # ---- COMPARE MODE ----
    else:
        st.markdown(
            '<div style="font-size:13px;color:#666;margin-bottom:16px;">'
            'Select two models to compare their defect profiles side by side.'
            '</div>',
            unsafe_allow_html=True
        )

        col_a, col_vs, col_b = st.columns([5, 1, 5])

        with col_a:
            st.markdown(
                '<div style="color:#E63946;font-family:IBM Plex Mono,'
                'monospace;font-size:11px;letter-spacing:1px;margin-bottom:8px;">'
                'MODEL A</div>',
                unsafe_allow_html=True
            )
            brand_a = st.selectbox(
                "Brand A", [""] + list(INDIAN_CAR_MODELS.keys()),
                key="brand_a"
            )
            model_opts_a = (
                [""] + INDIAN_CAR_MODELS.get(brand_a, [])
                if brand_a else [""]
            )
            model_a = st.selectbox("Model A", model_opts_a, key="model_a")
            year_a = st.selectbox("Year A", list(range(2025, 2017, -1)),
                                   key="year_a")

        with col_vs:
            st.markdown(
                '<div style="font-family:DM Serif Display,serif;'
                'font-size:32px;color:#E63946;text-align:center;'
                'padding-top:48px;">VS</div>',
                unsafe_allow_html=True
            )

        with col_b:
            st.markdown(
                '<div style="color:#2EC4B6;font-family:IBM Plex Mono,'
                'monospace;font-size:11px;letter-spacing:1px;margin-bottom:8px;">'
                'MODEL B</div>',
                unsafe_allow_html=True
            )
            brand_b = st.selectbox(
                "Brand B", [""] + list(INDIAN_CAR_MODELS.keys()),
                key="brand_b"
            )
            model_opts_b = (
                [""] + INDIAN_CAR_MODELS.get(brand_b, [])
                if brand_b else [""]
            )
            model_b = st.selectbox("Model B", model_opts_b, key="model_b")
            year_b = st.selectbox("Year B", list(range(2025, 2017, -1)),
                                   key="year_b")

        compare_btn = st.button(
            "⚖️ Compare Both Models",
            use_container_width=True,
            key="compare_btn"
        )

        if compare_btn and brand_a and model_a and brand_b and model_b:
            with st.status("Analysing both models...", expanded=True) as status:
                data_a, cached_a = run_analysis(
                    brand_a, model_a, year_a, secrets
                )
                data_b, cached_b = run_analysis(
                    brand_b, model_b, year_b, secrets
                )
                status.update(label="✅ Both analyses complete!",
                              state="complete")

            if data_a and data_b:
                res_a = data_a["result"]
                res_b = data_b["result"]
                esc_a = data_a["escalations"]
                esc_b = data_b["escalations"]

                # Find shared defects (simple title word overlap)
                titles_a = set(
                    " ".join(d.get("title","").lower().split()[:3])
                    for d in res_a.get("top_10_defects", [])
                )
                titles_b = set(
                    " ".join(d.get("title","").lower().split()[:3])
                    for d in res_b.get("top_10_defects", [])
                )
                shared_keys = titles_a & titles_b

                if shared_keys:
                    st.markdown(
                        '<div class="warning-box" style="border-color:#2EC4B6;'
                        'color:#2EC4B6;">🔗 Shared defects found across both '
                        f'models: {", ".join(shared_keys)}</div>',
                        unsafe_allow_html=True
                    )

                col_left, col_right = st.columns(2)

                with col_left:
                    st.markdown(
                        f'<div style="color:#E63946;font-family:DM Serif '
                        f'Display,serif;font-size:20px;margin-bottom:12px;">'
                        f'{brand_a} {model_a} {year_a}</div>',
                        unsafe_allow_html=True
                    )
                    if esc_a:
                        render_escalations(esc_a)
                    for i, d in enumerate(
                        res_a.get("top_10_defects", [])
                    ):
                        render_defect_card(d, i)

                with col_right:
                    st.markdown(
                        f'<div style="color:#2EC4B6;font-family:DM Serif '
                        f'Display,serif;font-size:20px;margin-bottom:12px;">'
                        f'{brand_b} {model_b} {year_b}</div>',
                        unsafe_allow_html=True
                    )
                    if esc_b:
                        render_escalations(esc_b)
                    for i, d in enumerate(
                        res_b.get("top_10_defects", [])
                    ):
                        render_defect_card(d, i)

        elif compare_btn:
            st.warning("Please select both models to compare.")


if __name__ == "__main__":
    main()
